"""
report_generator.py — Generate PPTX using the Lilly IBU template.

Uses the France reference deck as a template to inherit:
  - Branded slide masters (gradient headers, Lilly red)
  - Cover slide with hexagon graphics
  - Agenda layout with photo background
  - Content slide layout with footer + slide numbers
  - Thank You slide layout

Replaces data content while keeping branding pixel-perfect.
"""

from __future__ import annotations

import copy
import logging
import re
from datetime import datetime
from pathlib import Path

import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from lxml import etree

logger = logging.getLogger(__name__)

# Lilly colors
LILLY_RED = RGBColor(0xD5, 0x2B, 0x1E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x6B, 0x72, 0x80)
GREEN_FILL = RGBColor(0x6A, 0xB0, 0x4A)
PINK_FILL = RGBColor(0xF4, 0xC2, 0xC2)
GRAY_FILL = RGBColor(0xD9, 0xD9, 0xD9)
LIGHT_GREEN = RGBColor(0xC6, 0xEF, 0xCE)
LIGHT_YELLOW = RGBColor(0xFF, 0xF2, 0xCC)
LIGHT_RED = RGBColor(0xFF, 0xCC, 0xCC)
BLUE_FILL = RGBColor(0xB0, 0xC4, 0xDE)
TABLE_HDR = RGBColor(0xD5, 0x2B, 0x1E)
TABLE_ALT = RGBColor(0xFA, 0xF0, 0xF0)

TEMPLATE_PATH = None  # Set in config.yaml under output.template_pptx

def _fp(v):
    return f"{v:.0%}" if not pd.isna(v) else "N/A"

def _fn(v):
    return f"{int(v):,}" if not pd.isna(v) else "N/A"

def _set_cell(cell, text, sz=9, bold=False, color=DARK, align=PP_ALIGN.CENTER, fill=None):
    cell.text = str(text)
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(sz)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.alignment = align
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill

def _add_text(slide, left, top, width, height, text, sz=12, bold=False, color=DARK,
              align=PP_ALIGN.LEFT, name="Arial", italic=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = name
    p.font.italic = italic
    p.alignment = align
    return txBox

def _delete_all_slides(prs):
    """Delete all slides from the presentation, keeping masters/layouts."""
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        if rId is None:
            rId_attrib = list(prs.slides._sldIdLst[0].attrib.keys())
            for k in rId_attrib:
                if 'id' in k.lower() and 'r' in k.lower():
                    rId = prs.slides._sldIdLst[0].get(k)
                    break
        prs.part.drop_rel(rId)
        elem = prs.slides._sldIdLst[0]
        prs.slides._sldIdLst.remove(elem)

def _get_layout(prs, name):
    """Get layout by name. Falls back to blank layout if not found."""
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    # Fallback: try to find a blank or generic layout
    for layout in prs.slide_layouts:
        if "blank" in layout.name.lower() or layout.name == "Blank":
            return layout
    # Last resort: use the last layout
    return prs.slide_layouts[-1]


# ============================================================================
# Slide builders
# ============================================================================

def _slide_cover(prs, config):
    """Slide 1: Cover — uses 'Cover Slide 17' layout (red bg with hexagons)."""
    layout = _get_layout(prs, "Cover Slide 17")
    slide = prs.slides.add_slide(layout)

    start_raw = config.get("s3_parquet", {}).get("start_date", "2025-01-01")
    end_raw = config.get("s3_parquet", {}).get("end_date", "2025-12-31")
    # Format dates as "Jan 2025 – Dec 2025"
    try:
        from datetime import datetime as dt
        start_fmt = dt.strptime(start_raw, "%Y-%m-%d").strftime("%b %Y")
        end_fmt = dt.strptime(end_raw, "%Y-%m-%d").strftime("%b %Y")
    except:
        start_fmt, end_fmt = start_raw, end_raw
    gen_date = datetime.now().strftime("%B %d")

    country_code = config.get("_country_filter", None)
    country_name = config.get("_country_name", None)

    if country_name:
        title_text = f"IBU NBA/E UC Performance\nAnalysis & Rationalization ({country_name})"
    else:
        title_text = "IBU NBA/E UC Performance\nAnalysis & Rationalization"

    # Fill placeholders
    for ph in slide.placeholders:
        if "Title" in ph.name:
            ph.text = title_text
            ph.left = Inches(0.67)
            ph.top = Inches(2.28)
            ph.width = Inches(7.5)
            ph.height = Inches(1.5)
            for p in ph.text_frame.paragraphs:
                p.font.size = Pt(32)
                p.font.bold = True
                p.font.color.rgb = WHITE
                p.font.name = "Arial"
        elif "Text" in ph.name:
            ph.text = f"Assessment Period: {start_fmt} – {end_fmt}\n{gen_date}"
            ph.left = Inches(0.67)
            ph.top = Inches(4.5)
            ph.width = Inches(7.5)
            ph.height = Inches(1.0)
            for p in ph.text_frame.paragraphs:
                p.font.size = Pt(20)
                p.font.color.rgb = WHITE
                p.font.name = "Arial"


def _slide_objective(prs, config):
    """Slide 2: Objective — uses 'Cover Slide 3' layout."""
    layout = _get_layout(prs, "Cover Slide 3")
    slide = prs.slides.add_slide(layout)

    country_name = config.get("_country_name", None)
    if country_name:
        scope_text = f"in {country_name}"
    else:
        scope_text = "across all IBU markets (DE, ES, FR, GB, IT, PL)"

    for ph in slide.placeholders:
        if "Title" in ph.name:
            ph.text = "Objective"
            for p in ph.text_frame.paragraphs:
                p.font.size = Pt(36)
                p.font.bold = True
                p.font.color.rgb = WHITE
                p.font.name = "Arial"
        elif "Text" in ph.name:
            ph.text = (
                f"To assess the impact and effectiveness of IBU NBA/E Veeva suggestions "
                f"by measuring rep engagement and action rates {scope_text}. "
                f"This automated analysis identifies performance trends, anomalies, and "
                f"areas requiring attention to optimize omnichannel customer engagement."
            )
            for p in ph.text_frame.paragraphs:
                p.font.size = Pt(20)
                p.font.color.rgb = WHITE
                p.font.name = "Arial"


def _slide_executive_view(prs, df, kpi_results):
    """Slide 3: Executive View — uses layout engine to fill available space."""
    from src.slide_layout import Block, Gap, compute_layout

    layout = _get_layout(prs, "Agenda 4")
    slide = prs.slides.add_slide(layout)

    for ph in slide.placeholders:
        if "Title" in ph.name:
            ph.text = "Executive View"
            for p in ph.text_frame.paragraphs:
                p.font.size = Pt(36)
                p.font.bold = True
                p.font.color.rgb = WHITE
                p.font.name = "Arial"
        elif "Footer" in ph.name:
            ph.text = f"Company Confidential  © {datetime.now().year} Eli Lilly and Company"
        elif "Slide Number" in ph.name:
            ph.text = "3"
        elif "Picture" in ph.name:
            sp = ph._element
            sp.getparent().remove(sp)

    if len(df) == 0 or "hcp_id" not in df.columns:
        return

    # --- Data preparation ---
    RX = 4.5   # Right content area left edge
    RW = 8.5   # Right content area width

    df_w = df.copy()
    df_w["_cust_type"] = df_w["hcp_id"].apply(
        lambda x: "HCO" if str(x).endswith("HCO") else "HCP"
    )
    df_hcp = df_w[df_w["_cust_type"] == "HCP"]
    df_hco = df_w[df_w["_cust_type"] == "HCO"]
    has_hco = len(df_hco) > 0

    def _calc(data):
        t = len(data)
        c = data["hcp_id"].dropna().nunique() if t > 0 else 0
        rc = "rep_id" if "rep_id" in data.columns else "rep_name"
        r = data[rc].dropna().nunique() if rc in data.columns and t > 0 else 0
        return t, c, r, (int(t/c) if c > 0 else 0), (int(t/r) if r > 0 else 0)

    # Salesforce data
    sf_df = df_w.dropna(subset=["salesforce_sk", "medicine"]).copy()
    sf_data = []
    if len(sf_df) > 0:
        rc = "rep_id" if "rep_id" in sf_df.columns else "rep_name"
        sf = sf_df.groupby(["salesforce_sk", "_cust_type", "medicine"]).agg(
            suggestions=(sf_df.columns[0], "count"),
            reps=(rc, "nunique"),
            customers=("hcp_id", "nunique"),
        ).reset_index().sort_values("suggestions", ascending=False).head(10)
        sf_data = sf.to_dict("records")
    n_rows = len(sf_data)

    # --- Build layout blocks ---
    blocks = []
    blocks.append(Block("hcp_title", min_h=0.25, preferred_h=0.30, stretch=0.3))
    blocks.append(Block("hcp_cards", min_h=0.75, preferred_h=1.10, stretch=1.5))
    blocks.append(Gap(min_h=0.08, preferred_h=0.18, stretch=0.5))

    if has_hco:
        blocks.append(Block("hco_title", min_h=0.25, preferred_h=0.30, stretch=0.3))
        blocks.append(Block("hco_cards", min_h=0.75, preferred_h=1.10, stretch=1.5))
        blocks.append(Gap(min_h=0.08, preferred_h=0.18, stretch=0.5))

    blocks.append(Block("table", min_h=0.25 * (n_rows + 1), preferred_h=0.35 * (n_rows + 1), stretch=2.0))
    blocks.append(Gap(min_h=0.06, preferred_h=0.12, stretch=0.3))
    blocks.append(Block("note", min_h=0.65, preferred_h=0.90, stretch=1.0))

    compute_layout(blocks, top=0.08, bottom=6.48)

    # --- Helper: find block by name ---
    def _blk(name):
        for b in blocks:
            if hasattr(b, "name") and b.name == name:
                return b
        return None

    # --- Draw KPI cards ---
    def _draw_cards(blk_title, blk_cards, data, label):
        bt = _blk(blk_title)
        bc = _blk(blk_cards)
        if not bt or not bc:
            return

        _add_text(slide, RX, bt.y, RW, bt.actual_h,
                  label, 11, bold=True, color=DARK, align=PP_ALIGN.CENTER)

        t, c, r, ac, ar = _calc(data)
        vals = [
            ("Total\nSuggestions", _fn(t)),
            ("Total Customers\n(Reach)", _fn(c)),
            ("Total Reps", _fn(r)),
            ("Avg Sugg. / Cust.\n(Yearly)", str(ac)),
            ("Avg Sugg. / Rep\n(Yearly)", str(ar)),
        ]
        card_w = 1.55
        gap_w = (RW - 5 * card_w) / 4
        val_font = int(min(24, max(14, bc.actual_h * 18)))
        lbl_font = int(min(9, max(6, bc.actual_h * 7)))

        for i, (lbl, v) in enumerate(vals):
            x = RX + i * (card_w + gap_w)
            sh = slide.shapes.add_shape(5, Inches(x), Inches(bc.y), Inches(card_w), Inches(bc.actual_h))
            sh.fill.solid()
            sh.fill.fore_color.rgb = LILLY_RED
            sh.line.fill.background()
            tf = sh.text_frame
            tf.word_wrap = True
            p1 = tf.paragraphs[0]
            p1.text = lbl
            p1.font.size = Pt(lbl_font)
            p1.font.color.rgb = WHITE
            p1.font.name = "Arial"
            p1.alignment = PP_ALIGN.CENTER
            p2 = tf.add_paragraph()
            p2.text = v
            p2.font.size = Pt(val_font)
            p2.font.bold = True
            p2.font.color.rgb = WHITE
            p2.font.name = "Arial"
            p2.alignment = PP_ALIGN.CENTER

    _draw_cards("hcp_title", "hcp_cards", df_hcp, "HCP Summary (Actionable Suggestions)")
    if has_hco:
        _draw_cards("hco_title", "hco_cards", df_hco, "HCO Summary (Actionable Suggestions)")

    # --- Draw table ---
    bt = _blk("table")
    if bt and n_rows > 0:
        hdrs = ["Salesforces", "Customer", "Product", "Suggestions", "Reps",
                "Customers\nReached", "Avg Sugg./Cust.\n(Yearly)", "Avg Sugg./Rep\n(Yearly)"]
        nr = n_rows + 1
        row_h = bt.actual_h / nr
        tbl_font = int(min(9, max(7, row_h * 26)))

        tbl = slide.shapes.add_table(nr, 8, Inches(RX), Inches(bt.y), Inches(RW), Inches(bt.actual_h))
        table = tbl.table
        cw = [0.85, 0.7, 1.6, 0.95, 0.6, 0.95, 0.95, 0.9]
        for j, w in enumerate(cw):
            table.columns[j].width = Inches(w)
        for j, h in enumerate(hdrs):
            _set_cell(table.cell(0, j), h, tbl_font, bold=True, color=WHITE, fill=TABLE_HDR)
        for i, row in enumerate(sf_data):
            fc = TABLE_ALT if i % 2 == 0 else WHITE
            sg = row.get("suggestions", 0)
            cu = row.get("customers", 0)
            rp = row.get("reps", 0)
            sf_short = str(row.get("salesforce_sk", "")).replace("LLY3_CMRCL_", "")
            _set_cell(table.cell(i+1, 0), sf_short, tbl_font, fill=fc)
            _set_cell(table.cell(i+1, 1), str(row.get("_cust_type", "")), tbl_font, fill=fc)
            _set_cell(table.cell(i+1, 2), str(row.get("medicine", "")), tbl_font, fill=fc)
            _set_cell(table.cell(i+1, 3), _fn(sg), tbl_font, fill=fc)
            _set_cell(table.cell(i+1, 4), _fn(rp), tbl_font, fill=fc)
            _set_cell(table.cell(i+1, 5), _fn(cu), tbl_font, fill=fc)
            _set_cell(table.cell(i+1, 6), str(int(sg/cu)) if cu > 0 else "0", tbl_font, fill=fc)
            _set_cell(table.cell(i+1, 7), str(int(sg/rp)) if rp > 0 else "0", tbl_font, fill=fc)

    # --- Draw NOTE ---
    bn = _blk("note")
    if bn:
        note_title_font = int(min(12, max(8, bn.actual_h * 10)))
        note_body_font = int(min(10, max(7, bn.actual_h * 8)))

        nl = slide.shapes.add_shape(1, Inches(RX), Inches(bn.y), Inches(0.05), Inches(bn.actual_h))
        nl.fill.solid()
        nl.fill.fore_color.rgb = LILLY_RED
        nl.line.fill.background()
        nb = slide.shapes.add_textbox(Inches(RX + 0.15), Inches(bn.y), Inches(RW - 0.2), Inches(bn.actual_h))
        tf = nb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "NOTE"
        p.font.size = Pt(note_title_font)
        p.font.bold = True
        p.font.color.rgb = LILLY_RED
        p.font.name = "Arial"
        for note in [
            "Avg Suggestions per Customer = Total Suggestions / Total Customers (Reach)",
            "Avg Suggestions per Rep = Total Suggestions / Total Reps",
            "Excludes Non-Actionable Insights",
            "Sugg. = Suggestions & Cust. = Customers",
        ]:
            pn = tf.add_paragraph()
            pn.text = f"• {note}"
            pn.font.size = Pt(note_body_font)
            pn.font.color.rgb = DARK
            pn.font.name = "Arial"


def _slide_kpi_performance(prs, kpi_results, slide_num="4"):
    """Slide 4: Adherence Rate by Product — single red bar per product."""
    from src.slide_layout import Block, Gap, compute_layout

    layout = _get_layout(prs, "1_Title and Content")
    slide = prs.slides.add_slide(layout)

    for ph in slide.placeholders:
        if "Title" in ph.name:
            ph.text = "Product-level Performance"
        elif "Footer" in ph.name:
            ph.text = f"Company Confidential  © {datetime.now().year} Eli Lilly and Company"
        elif "Slide Number" in ph.name:
            ph.text = slide_num
        elif "Content" in ph.name:
            sp = ph._element
            sp.getparent().remove(sp)

    by_brand = kpi_results.get("by_brand", pd.DataFrame())
    if by_brand.empty:
        return

    by_brand = by_brand.sort_values("total_suggestions", ascending=False)
    n_brands = len(by_brand)

    # Chart title
    _add_text(slide, 0.5, 1.3, 9, 0.35,
              "Adherence Rate by Product (Adhered/Total Suggestions)", 14,
              bold=True, color=DARK, align=PP_ALIGN.CENTER)

    # Dynamic bar sizing — cap spacing so bars don't spread too far apart
    chart_top = 1.8
    chart_bottom = 5.5
    avail_h = chart_bottom - chart_top
    max_spacing = 0.85   # Never space bars more than this
    bar_spacing = min(max_spacing, avail_h / max(n_brands, 1))
    bar_h = min(0.42, bar_spacing * 0.6)

    # Center bars vertically in the chart area
    total_bars_h = n_brands * bar_spacing
    y_offset = chart_top + (avail_h - total_bars_h) / 2

    bar_left = 2.5      # Where bars start (after product name)
    bar_max_w = 7.5      # Maximum bar width (100%)

    # "Total Suggestions" header — just above the first bar
    _add_text(slide, 10.3, y_offset - 0.35, 2.2, 0.3, "Total Suggestions", 11,
              bold=True, color=LILLY_RED, align=PP_ALIGN.CENTER)

    # Draw bars
    for i, (_, row) in enumerate(by_brand.iterrows()):
        y = y_offset + i * bar_spacing
        med = str(row["medicine"])
        adh = float(row.get("adherence_rate", 0))
        total = int(row.get("total_suggestions", 0))

        # Product name
        _add_text(slide, 0.3, y, 2.1, bar_h, med, 11,
                  bold=True, color=DARK, align=PP_ALIGN.RIGHT)

        # Adherence bar (red)
        bar_w = max(0.3, bar_max_w * adh)
        sh = slide.shapes.add_shape(
            1, Inches(bar_left), Inches(y + 0.02), Inches(bar_w), Inches(bar_h - 0.04)
        )
        sh.fill.solid()
        sh.fill.fore_color.rgb = LILLY_RED
        sh.line.fill.background()

        # Percentage label (at end of bar)
        _add_text(slide, bar_left + bar_w + 0.1, y, 0.6, bar_h,
                  _fp(adh), 11, bold=True, color=DARK, align=PP_ALIGN.LEFT)

        # Total suggestions (right side, in rounded box)
        box = slide.shapes.add_shape(
            5, Inches(10.5), Inches(y), Inches(1.8), Inches(bar_h)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xE8, 0xEA, 0xF6)
        box.line.color.rgb = RGBColor(0xB0, 0xBE, 0xD0)
        box.line.width = Pt(0.5)
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = _fn(total)
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = LILLY_RED
        p.font.name = "Arial"
        p.alignment = PP_ALIGN.CENTER

    # X-axis position (below the last bar)
    axis_y = y_offset + n_brands * bar_spacing + 0.1

    # Y-axis line (vertical, left edge of bars)
    y_axis_top = y_offset - 0.05
    y_axis_h = (n_brands * bar_spacing) + 0.15
    y_line = slide.shapes.add_shape(
        1, Inches(bar_left - 0.02), Inches(y_axis_top), Inches(0.015), Inches(y_axis_h)
    )
    y_line.fill.solid()
    y_line.fill.fore_color.rgb = MED_GRAY
    y_line.line.fill.background()

    # X-axis line
    line = slide.shapes.add_shape(
        1, Inches(bar_left), Inches(axis_y), Inches(bar_max_w), Inches(0.015)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = MED_GRAY
    line.line.fill.background()

    # X-axis labels (0% to 100%)
    for pct in range(0, 110, 10):
        x = bar_left + bar_max_w * (pct / 100)
        _add_text(slide, x - 0.3, axis_y + 0.05, 0.6, 0.2,
                  f"{pct}%", 9, color=MED_GRAY, align=PP_ALIGN.CENTER)


def _slide_product_detail(prs, kpi_results, slide_num="5"):
    """Slide 5: Acceptance vs Dismissal + Dismissal reasons — dynamically sized."""
    layout = _get_layout(prs, "1_Title and Content")
    slide = prs.slides.add_slide(layout)

    for ph in slide.placeholders:
        if "Title" in ph.name:
            ph.text = "Product-level Performance"
        elif "Footer" in ph.name:
            ph.text = f"Company Confidential  © {datetime.now().year} Eli Lilly and Company"
        elif "Slide Number" in ph.name:
            ph.text = slide_num
        elif "Content" in ph.name:
            sp = ph._element
            sp.getparent().remove(sp)

    by_brand = kpi_results.get("by_brand", pd.DataFrame())
    if by_brand.empty:
        return
    by_brand = by_brand.sort_values("total_suggestions", ascending=False)
    n_brands = len(by_brand)

    dismissals = kpi_results.get("dismissal_reasons", pd.DataFrame())
    reason_col = "reason" if not dismissals.empty and "reason" in dismissals.columns else "reason_clean"
    top_reasons = dismissals.head(5) if not dismissals.empty else pd.DataFrame()
    n_reasons = len(top_reasons)

    # --- LEFT SIDE: Acceptance vs Dismissal within Adherence ---
    left_x = 0.3
    left_w = 6.2

    # Dynamic bar sizing
    chart_top = 2.05
    chart_bottom = 5.2
    avail_h = chart_bottom - chart_top
    max_spacing = 0.85
    bar_spacing = min(max_spacing, avail_h / max(n_brands, 1))
    bar_h = min(0.42, bar_spacing * 0.6)
    total_bars = n_brands * bar_spacing
    y_start = chart_top + (avail_h - total_bars) / 2

    name_w = 1.6
    bar_left = left_x + name_w + 0.1
    bar_max_w = 3.5

    # Left chart title — centered below header bar
    _add_text(slide, left_x, 1.45, left_w, 0.3,
              "Acceptance vs Dismissal within Adherence by Product", 10,
              bold=True, color=DARK, align=PP_ALIGN.CENTER)

    # "Adhered Suggestions" header — above first bar
    _add_text(slide, bar_left + bar_max_w + 0.15, y_start - 0.35, 1.0, 0.3,
              "Adhered\nSuggestions", 8, bold=True, color=LILLY_RED, align=PP_ALIGN.CENTER)

    for i, (_, row) in enumerate(by_brand.iterrows()):
        y = y_start + i * bar_spacing
        med = str(row["medicine"])
        acc_r = float(row.get("adherence_due_to_acceptance", 0))
        dis_r = float(row.get("adherence_due_to_dismissal", 0))
        adhered = int(row.get("accept_count", 0)) + int(row.get("dismiss_count", 0))

        # Product name
        _add_text(slide, left_x, y, name_w, bar_h, med, 9,
                  bold=True, color=DARK, align=PP_ALIGN.RIGHT)

        # Green bar (acceptance portion)
        acc_w = max(0.01, bar_max_w * acc_r)
        sh1 = slide.shapes.add_shape(1, Inches(bar_left), Inches(y + 0.02), Inches(acc_w), Inches(bar_h - 0.04))
        sh1.fill.solid()
        sh1.fill.fore_color.rgb = GREEN_FILL
        sh1.line.fill.background()
        if acc_r > 0.08:
            _add_text(slide, bar_left, y + 0.02, acc_w, bar_h - 0.04,
                      _fp(acc_r), 8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Red bar (dismissal portion)
        dis_w = max(0.01, bar_max_w * dis_r) if dis_r > 0.01 else 0
        if dis_w > 0:
            sh2 = slide.shapes.add_shape(1, Inches(bar_left + acc_w), Inches(y + 0.02), Inches(dis_w), Inches(bar_h - 0.04))
            sh2.fill.solid()
            sh2.fill.fore_color.rgb = LILLY_RED
            sh2.line.fill.background()
            if dis_r > 0.08:
                _add_text(slide, bar_left + acc_w, y + 0.02, dis_w, bar_h - 0.04,
                          _fp(dis_r), 8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Adhered count
        _add_text(slide, bar_left + bar_max_w + 0.15, y, 1.0, bar_h,
                  _fn(adhered), 9, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    # Left Y-axis line
    y_axis_l = slide.shapes.add_shape(
        1, Inches(bar_left - 0.02), Inches(y_start - 0.05),
        Inches(0.015), Inches(n_brands * bar_spacing + 0.13)
    )
    y_axis_l.fill.solid()
    y_axis_l.fill.fore_color.rgb = MED_GRAY
    y_axis_l.line.fill.background()

    # Left X-axis
    axis_y_l = y_start + n_brands * bar_spacing + 0.08
    line_l = slide.shapes.add_shape(1, Inches(bar_left), Inches(axis_y_l), Inches(bar_max_w), Inches(0.012))
    line_l.fill.solid()
    line_l.fill.fore_color.rgb = MED_GRAY
    line_l.line.fill.background()
    for pct in [0, 20, 40, 60, 80, 100]:
        x = bar_left + bar_max_w * (pct / 100)
        _add_text(slide, x - 0.25, axis_y_l + 0.03, 0.5, 0.18, f"{pct}%", 8, color=MED_GRAY, align=PP_ALIGN.CENTER)

    # Left legend
    leg_y = axis_y_l + 0.25
    sh_g = slide.shapes.add_shape(1, Inches(bar_left + 0.5), Inches(leg_y), Inches(0.2), Inches(0.15))
    sh_g.fill.solid()
    sh_g.fill.fore_color.rgb = GREEN_FILL
    sh_g.line.fill.background()
    _add_text(slide, bar_left + 0.75, leg_y - 0.02, 1.0, 0.18, "Acceptance", 8, color=DARK)
    sh_r = slide.shapes.add_shape(1, Inches(bar_left + 2.0), Inches(leg_y), Inches(0.2), Inches(0.15))
    sh_r.fill.solid()
    sh_r.fill.fore_color.rgb = LILLY_RED
    sh_r.line.fill.background()
    _add_text(slide, bar_left + 2.25, leg_y - 0.02, 1.0, 0.18, "Dismissal", 8, color=DARK)

    # --- RIGHT SIDE: Dismissal Reasons ---
    right_x = 6.8
    right_w = 6.0
    if n_reasons > 0:
        _add_text(slide, right_x, 1.45, right_w, 0.3,
                  "Share of Dismissal Reasons \u2013 All Products (Top 5)", 10,
                  bold=True, color=DARK, align=PP_ALIGN.CENTER)

        r_name_w = 2.8
        r_bar_left = right_x + r_name_w + 0.1
        r_bar_max_w = 2.0

        # Dynamic sizing for reasons
        r_spacing = min(max_spacing, avail_h / max(n_reasons, 1))
        r_bar_h = min(0.42, r_spacing * 0.55)
        r_total = n_reasons * r_spacing
        r_y_start = chart_top + (avail_h - r_total) / 2

        # Find max pct for axis scaling
        max_pct = float(top_reasons["pct"].max()) if not top_reasons.empty else 0.5
        axis_max = min(1.0, max(0.3, (int(max_pct * 10) + 2) / 10))

        for i, (_, row) in enumerate(top_reasons.iterrows()):
            y = r_y_start + i * r_spacing
            reason = str(row.get(reason_col, ""))
            if len(reason) > 50:
                reason = reason[:47] + "..."
            pct_val = float(row.get("pct", 0))

            # Reason text
            _add_text(slide, right_x, y, r_name_w, r_bar_h, reason, 8, color=DARK, align=PP_ALIGN.RIGHT)

            # Red bar
            r_w = max(0.15, r_bar_max_w * (pct_val / axis_max))
            sh = slide.shapes.add_shape(1, Inches(r_bar_left), Inches(y + 0.03), Inches(r_w), Inches(r_bar_h - 0.06))
            sh.fill.solid()
            sh.fill.fore_color.rgb = LILLY_RED
            sh.line.fill.background()

            # Percentage label
            _add_text(slide, r_bar_left + r_w + 0.05, y, 0.7, r_bar_h,
                      f"{pct_val:.1%}", 8, bold=True, color=DARK, align=PP_ALIGN.LEFT)

        # Right Y-axis line
        y_axis_r = slide.shapes.add_shape(
            1, Inches(r_bar_left - 0.02), Inches(r_y_start - 0.05),
            Inches(0.015), Inches(n_reasons * r_spacing + 0.13)
        )
        y_axis_r.fill.solid()
        y_axis_r.fill.fore_color.rgb = MED_GRAY
        y_axis_r.line.fill.background()

        # Right X-axis
        axis_y_r = r_y_start + n_reasons * r_spacing + 0.08
        line_r = slide.shapes.add_shape(1, Inches(r_bar_left), Inches(axis_y_r), Inches(r_bar_max_w), Inches(0.012))
        line_r.fill.solid()
        line_r.fill.fore_color.rgb = MED_GRAY
        line_r.line.fill.background()
        step = max(0.1, round(axis_max / 5, 1))
        p = 0.0
        while p <= axis_max + 0.001:
            x = r_bar_left + r_bar_max_w * (p / axis_max)
            _add_text(slide, x - 0.25, axis_y_r + 0.03, 0.5, 0.18, f"{p:.0%}", 8, color=MED_GRAY, align=PP_ALIGN.CENTER)
            p += step

        # Right legend
        sh_rl = slide.shapes.add_shape(1, Inches(r_bar_left + 0.3), Inches(axis_y_r + 0.25), Inches(0.2), Inches(0.15))
        sh_rl.fill.solid()
        sh_rl.fill.fore_color.rgb = LILLY_RED
        sh_rl.line.fill.background()
        _add_text(slide, r_bar_left + 0.55, axis_y_r + 0.23, 1.5, 0.18, "Dismissal Reasons", 8, color=DARK)

    # --- NOTE section (bottom left) ---
    note_y = 5.95
    nl = slide.shapes.add_shape(1, Inches(left_x), Inches(note_y), Inches(0.05), Inches(0.55))
    nl.fill.solid()
    nl.fill.fore_color.rgb = LILLY_RED
    nl.line.fill.background()
    nb = slide.shapes.add_textbox(Inches(left_x + 0.15), Inches(note_y), Inches(5.8), Inches(0.55))
    tf_n = nb.text_frame
    tf_n.word_wrap = True
    p_n = tf_n.paragraphs[0]
    p_n.text = "NOTE"
    p_n.font.size = Pt(8)
    p_n.font.bold = True
    p_n.font.color.rgb = LILLY_RED
    p_n.font.name = "Arial"
    for note in [
        "Excludes Non-Actionable Insights (Pinned Insight, Insight)",
        "Adherence = Accepted + Dismissed (reps who acted on the suggestion)",
    ]:
        pn = tf_n.add_paragraph()
        pn.text = f"• {note}"
        pn.font.size = Pt(7)
        pn.font.color.rgb = DARK
        pn.font.name = "Arial"

    # --- TOP 2 DISMISSAL REASONS box (bottom right) with red strip ---
    if n_reasons >= 2:
        box_y = 5.95
        # Red left strip
        tr = slide.shapes.add_shape(1, Inches(right_x), Inches(box_y), Inches(0.05), Inches(0.55))
        tr.fill.solid()
        tr.fill.fore_color.rgb = LILLY_RED
        tr.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(right_x + 0.15), Inches(box_y), Inches(right_w - 0.15), Inches(0.55))
        tf_t = tb.text_frame
        tf_t.word_wrap = True
        p_t = tf_t.paragraphs[0]
        p_t.text = "TOP 2 DISMISSAL REASONS"
        p_t.font.size = Pt(8)
        p_t.font.bold = True
        p_t.font.color.rgb = LILLY_RED
        p_t.font.name = "Arial"
        for j in range(min(2, n_reasons)):
            r = top_reasons.iloc[j]
            reason_text = str(r.get(reason_col, ""))[:60]
            pct_text = f"{r['pct']:.1%}" if 'pct' in r else ""
            pt = tf_t.add_paragraph()
            pt.text = f"{j+1}. {reason_text} ({pct_text})"
            pt.font.size = Pt(7)
            pt.font.color.rgb = DARK
            pt.font.name = "Arial"


def _slide_usecase_heatmap(prs, kpi_results, slide_num="6"):
    """Slide 6: Use Case Heatmap — green gradient matching reference deck."""
    layout = _get_layout(prs, "1_Title and Content")
    slide = prs.slides.add_slide(layout)

    for ph in slide.placeholders:
        if "Title" in ph.name:
            ph.text = "Use Case-level Performance"
        elif "Footer" in ph.name:
            ph.text = f"Company Confidential  © {datetime.now().year} Eli Lilly and Company"
        elif "Slide Number" in ph.name:
            ph.text = slide_num
        elif "Content" in ph.name:
            sp = ph._element
            sp.getparent().remove(sp)

    _add_text(slide, 0.67, 1.45, 12, 0.3,
              "Acceptance Rate (Total) - Product x Use Case Heatmap", 13,
              bold=True, color=DARK, align=PP_ALIGN.CENTER)

    by_bu = kpi_results.get("by_brand_usecase", pd.DataFrame())
    if by_bu.empty:
        return

    pivot = by_bu.pivot_table(
        index="medicine", columns="sugg_name",
        values="acceptance_rate_total", aggfunc="first"
    ).fillna(-1)

    uc_totals = by_bu.groupby("sugg_name")["total_suggestions"].sum().sort_values(ascending=False)
    top_ucs = uc_totals.head(10).index.tolist()
    pivot = pivot[[c for c in top_ucs if c in pivot.columns]]

    products = pivot.index.tolist()
    use_cases = pivot.columns.tolist()
    if not products or not use_cases:
        return

    # Green gradient colors matching the reference deck
    GREEN_HIGH = RGBColor(0x00, 0xB0, 0x50)
    GREEN_MED = RGBColor(0x92, 0xD0, 0x50)
    GREEN_LOW_MED = RGBColor(0xC6, 0xE0, 0xB4)
    GREEN_LOW = RGBColor(0xE2, 0xEF, 0xDA)

    def _hc(val):
        if val >= 0.71: return GREEN_HIGH
        elif val >= 0.51: return GREEN_MED
        elif val >= 0.31: return GREEN_LOW_MED
        else: return GREEN_LOW

    def _tc(val):
        return WHITE if val >= 0.71 else DARK

    # Dynamic sizing to fill the slide
    n_rows = len(products) + 1
    n_cols = len(use_cases) + 1

    # Available vertical space: 1.85 to 6.45
    avail_top = 1.85
    avail_bottom = 6.4
    avail_h = avail_bottom - avail_top

    # Allocate: table gets most space, note+legend at bottom
    note_h = 0.55
    gap = 0.2
    table_h = avail_h - note_h - gap
    row_h = table_h / n_rows
    row_h = min(0.55, max(0.32, row_h))  # Clamp
    actual_table_h = row_h * n_rows

    # Center table in available space above note
    table_top = avail_top + ((avail_h - note_h - gap - actual_table_h) / 2)
    col_w = min(1.3, 12.0 / n_cols)
    table_w = col_w * n_cols

    tbl = slide.shapes.add_table(n_rows, n_cols,
        Inches(0.67), Inches(table_top), Inches(table_w), Inches(actual_table_h))
    table = tbl.table

    font_sz = int(min(10, max(7, row_h * 22)))
    hdr_sz = int(min(9, max(6, row_h * 18)))

    _set_cell(table.cell(0, 0), "Product", font_sz, bold=True, color=WHITE, fill=TABLE_HDR, align=PP_ALIGN.LEFT)
    for j, uc in enumerate(use_cases):
        _set_cell(table.cell(0, j+1), str(uc)[:18], hdr_sz, bold=True, color=WHITE, fill=TABLE_HDR)

    for i, prod in enumerate(products):
        row_bg = TABLE_ALT if i % 2 == 0 else WHITE
        _set_cell(table.cell(i+1, 0), prod, font_sz, bold=True, color=DARK, fill=row_bg, align=PP_ALIGN.LEFT)
        for j, uc in enumerate(use_cases):
            val = pivot.loc[prod, uc]
            if val == -1 or pd.isna(val):
                _set_cell(table.cell(i+1, j+1), "", font_sz, fill=WHITE)
            else:
                _set_cell(table.cell(i+1, j+1), _fp(val), font_sz, color=_tc(val), fill=_hc(val))

    # NOTE + Legend positioned at bottom
    bottom_y = avail_bottom - note_h

    # Legend (bottom right)
    legend_tbl = slide.shapes.add_table(1, 3,
        Inches(7.0), Inches(bottom_y), Inches(5.66), Inches(0.4))
    lt = legend_tbl.table
    _set_cell(lt.cell(0, 0), "Low Acceptance\n(0% - 30%)", 8, color=DARK, fill=GREEN_LOW)
    _set_cell(lt.cell(0, 1), "Medium Acceptance\n(31% - 70%)", 8, color=DARK, fill=GREEN_MED)
    _set_cell(lt.cell(0, 2), "High Acceptance\n(71% - 100%)", 8, color=WHITE, fill=GREEN_HIGH)


def _slide_usecase_performance(prs, kpi_results, start_slide_num=7):
    """Slides 7+: Use Case per brand — matching reference layout exactly."""
    by_bu = kpi_results.get("by_brand_usecase", pd.DataFrame())
    if by_bu.empty:
        return start_slide_num

    brands = by_bu.groupby("medicine")["total_suggestions"].sum().sort_values(ascending=False).index.tolist()
    slide_num = start_slide_num

    MAX_UCS_PER_SLIDE = 8

    for brand in brands:
        brand_data = by_bu[by_bu["medicine"] == brand].sort_values("total_suggestions", ascending=False)
        if len(brand_data) == 0:
            continue

        chunks = [brand_data.iloc[i:i + MAX_UCS_PER_SLIDE]
                  for i in range(0, len(brand_data), MAX_UCS_PER_SLIDE)]

        for chunk_idx, chunk in enumerate(chunks):
            layout = _get_layout(prs, "1_Title and Content")
            slide = prs.slides.add_slide(layout)

            for ph in slide.placeholders:
                if "Title" in ph.name:
                    ph.text = "Use Case-level Performance"
                elif "Footer" in ph.name:
                    ph.text = f"Company Confidential  © {datetime.now().year} Eli Lilly and Company"
                elif "Slide Number" in ph.name:
                    ph.text = str(slide_num)
                elif "Content" in ph.name:
                    sp = ph._element
                    sp.getparent().remove(sp)

            n_ucs = len(chunk)

            # Brand name (top left, matching reference)
            brand_label = brand if chunk_idx == 0 else f"{brand} (cont.)"
            _add_text(slide, 0.67, 1.35, 2.5, 0.25, brand_label, 11,
                      bold=True, color=LILLY_RED, align=PP_ALIGN.LEFT)

            # Column headers (matching reference positions and colors)
            _add_text(slide, 3.0, 1.35, 2.3, 0.25, "Accepted", 10, bold=True,
                      color=GREEN_FILL, align=PP_ALIGN.CENTER)
            _add_text(slide, 5.5, 1.35, 2.3, 0.25, "Dismissed", 10, bold=True,
                      color=LILLY_RED, align=PP_ALIGN.CENTER)
            _add_text(slide, 7.8, 1.35, 1.5, 0.25, "No Action", 10, bold=True,
                      color=MED_GRAY, align=PP_ALIGN.CENTER)
            _add_text(slide, 10.0, 1.35, 2.5, 0.25, "Total Suggestions", 10, bold=True,
                      color=LILLY_RED, align=PP_ALIGN.CENTER)

            # Dynamic row sizing
            chart_top = 1.7
            chart_bottom = 4.7
            avail = chart_bottom - chart_top
            row_spacing = min(0.5, avail / max(n_ucs, 1))
            bar_h = min(0.32, row_spacing * 0.7)
            total_h = n_ucs * row_spacing
            y_start = chart_top + (avail - total_h) / 2

            # Bar widths (wider to fill the column space)
            acc_bar_w = 2.0      # Fixed width for all bars
            dis_bar_w = 2.0
            noa_bar_w = 2.0

            # Green shades based on acceptance rate
            def _acc_green(val):
                if val >= 0.80: return RGBColor(0x6A, 0xB0, 0x4A)
                elif val >= 0.60: return RGBColor(0x92, 0xCF, 0x50)
                elif val >= 0.40: return RGBColor(0xA9, 0xD1, 0x8E)
                else: return RGBColor(0xC6, 0xE0, 0xB4)

            # Red shades based on dismissal rate
            def _dis_red(val):
                if val >= 0.35: return RGBColor(0xD5, 0x2B, 0x1E)   # Dark red
                elif val >= 0.20: return RGBColor(0xE8, 0x5D, 0x50) # Medium red
                elif val >= 0.10: return RGBColor(0xF4, 0x8F, 0x8F) # Light red
                else: return RGBColor(0xF4, 0xC2, 0xC2)             # Pink

            # Gray shades based on no action rate
            def _noa_gray(val):
                if val >= 0.50: return RGBColor(0xA0, 0xA0, 0xA0)   # Dark gray
                elif val >= 0.30: return RGBColor(0xBF, 0xBF, 0xBF) # Medium gray
                else: return RGBColor(0xD9, 0xD9, 0xD9)             # Light gray

            for row_idx, (_, row) in enumerate(chunk.iterrows()):
                y = y_start + row_idx * row_spacing
                uc_name = str(row["sugg_name"])
                acc = float(row.get("acceptance_rate_total", 0))
                dis = float(row.get("dismissal_rate_total", 0))
                noa = float(row.get("no_action_rate", 0))
                total = int(row.get("total_suggestions", 0))

                # Use case name (indented under brand)
                _add_text(slide, 0.8, y, 2.0, bar_h, uc_name, 9, color=DARK, align=PP_ALIGN.LEFT)

                # Accepted bar (green, shade varies by value)
                acc_w = acc_bar_w
                sh = slide.shapes.add_shape(1, Inches(3.0), Inches(y + 0.01),
                                            Inches(acc_w), Inches(bar_h - 0.02))
                sh.fill.solid()
                sh.fill.fore_color.rgb = _acc_green(acc)
                sh.line.fill.background()
                _add_text(slide, 3.0, y + 0.01, acc_w, bar_h - 0.02,
                          _fp(acc), 8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

                # Dismissed bar (fixed width, shade varies)
                if dis > 0.01:
                    dis_w = dis_bar_w
                    sh2 = slide.shapes.add_shape(1, Inches(5.5), Inches(y + 0.01),
                                                 Inches(dis_w), Inches(bar_h - 0.02))
                    sh2.fill.solid()
                    sh2.fill.fore_color.rgb = _dis_red(dis)
                    sh2.line.fill.background()
                    _add_text(slide, 5.5, y + 0.01, dis_w, bar_h - 0.02,
                              _fp(dis), 8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
                else:
                    sh2 = slide.shapes.add_shape(1, Inches(5.5), Inches(y + 0.01),
                                                 Inches(dis_bar_w), Inches(bar_h - 0.02))
                    sh2.fill.solid()
                    sh2.fill.fore_color.rgb = RGBColor(0xFA, 0xE0, 0xE0)
                    sh2.line.fill.background()
                    _add_text(slide, 5.5, y + 0.01, dis_bar_w, bar_h - 0.02,
                              "0%", 8, bold=True, color=MED_GRAY, align=PP_ALIGN.CENTER)

                # No Action bar (fixed width, shade varies)
                if noa > 0.01:
                    noa_w = noa_bar_w
                    sh3 = slide.shapes.add_shape(1, Inches(7.8), Inches(y + 0.01),
                                                 Inches(noa_w), Inches(bar_h - 0.02))
                    sh3.fill.solid()
                    sh3.fill.fore_color.rgb = _noa_gray(noa)
                    sh3.line.fill.background()
                    _add_text(slide, 7.8, y + 0.01, noa_w, bar_h - 0.02,
                              _fp(noa), 8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
                else:
                    sh3 = slide.shapes.add_shape(1, Inches(7.8), Inches(y + 0.01),
                                                 Inches(noa_bar_w), Inches(bar_h - 0.02))
                    sh3.fill.solid()
                    sh3.fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
                    sh3.line.fill.background()
                    _add_text(slide, 7.8, y + 0.01, noa_bar_w, bar_h - 0.02,
                              "0%", 8, bold=True, color=MED_GRAY, align=PP_ALIGN.CENTER)

                # Total suggestions (blue rounded box)
                tb = slide.shapes.add_shape(5, Inches(10.2), Inches(y), Inches(1.8), Inches(bar_h))
                tb.fill.solid()
                tb.fill.fore_color.rgb = RGBColor(0xE8, 0xEA, 0xF6)
                tb.line.color.rgb = RGBColor(0xB0, 0xBE, 0xD0)
                tb.line.width = Pt(0.5)
                tf = tb.text_frame
                p = tf.paragraphs[0]
                p.text = _fn(total)
                p.font.size = Pt(9)
                p.font.bold = True
                p.font.color.rgb = LILLY_RED
                p.font.name = "Arial"
                p.alignment = PP_ALIGN.CENTER

            # --- Key Callouts table ---
            callouts = []
            for _, row in chunk.iterrows():
                uc = str(row["sugg_name"])
                acc_v = float(row.get("acceptance_rate_total", 0))
                dis_v = float(row.get("dismissal_rate_total", 0))
                noa_v = float(row.get("no_action_rate", 0))
                if acc_v < 0.50:
                    callouts.append((uc, f"{_fp(acc_v)} (Acceptance Rate)"))
                if dis_v > 0.20:
                    callouts.append((uc, f"{_fp(dis_v)} (Dismissal Rate)"))
                if noa_v > 0.20:
                    callouts.append((uc, f"{_fp(noa_v)} (No Action Rate)"))

            callout_y = 5.0
            _add_text(slide, 0.67, callout_y - 0.22, 12, 0.2,
                      "Key Callouts \u2014 Use Cases Below Threshold (<50% Acceptance or  >20% Dismissal/No Action Rate)",
                      8, bold=True, color=LILLY_RED)

            if callouts:
                n_co = min(len(callouts), 6)
                co_row_h = min(0.25, (6.3 - callout_y) / (n_co + 1))
                co_tbl = slide.shapes.add_table(n_co + 1, 2,
                    Inches(0.67), Inches(callout_y), Inches(12.0), Inches(co_row_h * (n_co + 1)))
                ct = co_tbl.table
                ct.columns[0].width = Inches(5.5)
                ct.columns[1].width = Inches(6.5)
                _set_cell(ct.cell(0, 0), "Use Case Performing Below Threshold", 7,
                          bold=True, color=WHITE, fill=TABLE_HDR)
                _set_cell(ct.cell(0, 1), "Key Callouts/ Potential Reasons", 7,
                          bold=True, color=WHITE, fill=TABLE_HDR)
                for ci in range(n_co):
                    uc, detail = callouts[ci]
                    _set_cell(ct.cell(ci + 1, 0), f"{uc} \u2014 {detail}", 7,
                              fill=WHITE, align=PP_ALIGN.LEFT)
                    _set_cell(ct.cell(ci + 1, 1), "", 7, fill=WHITE, align=PP_ALIGN.LEFT)
            else:
                _add_text(slide, 0.67, callout_y, 12, 0.25,
                          "All use cases performing within threshold.", 8, color=GREEN_FILL)

            slide_num += 1

    return slide_num


def _slide_anomalies(prs, anomalies_df, slide_num="7"):
    """Slide 7: Anomalies — table with severity coloring."""
    layout = _get_layout(prs, "1_Title and Content")
    slide = prs.slides.add_slide(layout)

    for ph in slide.placeholders:
        if "Title" in ph.name:
            ph.text = "Anomalies & Areas of Concern"
        elif "Footer" in ph.name:
            ph.text = f"Company Confidential  © {datetime.now().year} Eli Lilly and Company"
        elif "Slide Number" in ph.name:
            ph.text = slide_num
        elif "Content" in ph.name:
            sp = ph._element
            sp.getparent().remove(sp)

    if anomalies_df.empty:
        _add_text(slide, 2, 3, 9, 1,
                  "No anomalies detected — all metrics within normal range.", 20,
                  color=GREEN_FILL, align=PP_ALIGN.CENTER)
        return

    items = anomalies_df.head(12)
    headers = ["Severity", "Market", "Metric", "Details"]
    n_rows = len(items) + 1
    tbl = slide.shapes.add_table(n_rows, 4,
        Inches(0.67), Inches(1.5), Inches(12.0), Inches(0.38 * n_rows))
    table = tbl.table
    table.columns[0].width = Inches(1.2)
    table.columns[1].width = Inches(1.5)
    table.columns[2].width = Inches(2.5)
    table.columns[3].width = Inches(6.8)

    for j, h in enumerate(headers):
        _set_cell(table.cell(0, j), h, 9, bold=True, color=WHITE, fill=TABLE_HDR)

    for i, (_, row) in enumerate(items.iterrows()):
        sev = str(row["severity"]).upper()
        bg = LIGHT_RED if sev == "CRITICAL" else LIGHT_YELLOW
        sev_color = LILLY_RED if sev == "CRITICAL" else RGBColor(0xD9, 0x77, 0x06)

        _set_cell(table.cell(i+1, 0), sev, 8, bold=True, color=sev_color, fill=bg)
        _set_cell(table.cell(i+1, 1), str(row["dimension"]), 8, fill=bg, align=PP_ALIGN.LEFT)
        _set_cell(table.cell(i+1, 2), str(row["metric"]), 8, fill=bg, align=PP_ALIGN.LEFT)
        explanation = str(row["explanation"])
        if len(explanation) > 85:
            explanation = explanation[:82] + "..."
        _set_cell(table.cell(i+1, 3), explanation, 8, fill=bg, align=PP_ALIGN.LEFT)


def _slide_appendix_cover(prs):
    """Appendix cover — uses 'Cover Slide 3' layout (gradient background)."""
    layout = _get_layout(prs, "Cover Slide 3")
    slide = prs.slides.add_slide(layout)
    for ph in slide.placeholders:
        if "Title" in ph.name:
            ph.text = "Appendix"
            for p in ph.text_frame.paragraphs:
                p.font.size = Pt(36)
                p.font.bold = True
                p.font.color.rgb = WHITE
                p.font.name = "Arial"
                p.alignment = PP_ALIGN.CENTER
        elif "Text" in ph.name:
            ph.text = ""


def _slide_appendix_content(prs, config, slide_num):
    """Appendix content — Key Metrics Definitions + Scope matching reference."""
    layout = _get_layout(prs, "1_Title and Content")
    slide = prs.slides.add_slide(layout)

    for ph in slide.placeholders:
        if "Title" in ph.name:
            ph.text = "Appendix"
        elif "Footer" in ph.name:
            ph.text = f"Company Confidential  © {datetime.now().year} Eli Lilly and Company"
        elif "Slide Number" in ph.name:
            ph.text = str(slide_num)
        elif "Content" in ph.name:
            sp = ph._element
            sp.getparent().remove(sp)

    # --- Bordered box for KEY METRICS DEFINITIONS ---
    box1 = slide.shapes.add_shape(5, Inches(0.5), Inches(1.55), Inches(12.3), Inches(3.0))
    box1.fill.solid()
    box1.fill.fore_color.rgb = WHITE
    box1.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    box1.line.width = Pt(1)

    # Red header bar
    hdr1 = slide.shapes.add_shape(5, Inches(1.2), Inches(1.7), Inches(3.5), Inches(0.32))
    hdr1.fill.solid()
    hdr1.fill.fore_color.rgb = LILLY_RED
    hdr1.line.fill.background()
    tf = hdr1.text_frame
    p = tf.paragraphs[0]
    p.text = "KEY METRICS DEFINITIONS"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Arial"

    # Metric definitions with numbered red circles
    definitions = [
        "Adherence Rate = { (Accepted Suggestion + Dismissed Suggestion) / Total Actionable Suggestion } x 100",
        "Acceptance Rate = (Accepted Suggestion / Total Actionable Suggestion) × 100",
        "Dismissed Rate = (Dismissed Suggestion / Total Actionable Suggestion) × 100",
        "No Action Rate = (Ignored Suggestion / Total Actionable Suggestion) × 100",
        "Avg Sugg. per Customer (Yearly) = Total Actionable Suggestions / Total Customers",
        "Avg Sugg. per Rep (Yearly) = Total Actionable Suggestions / Total Reps",
    ]
    for i, defn in enumerate(definitions):
        y = 2.25 + i * 0.35
        # Red circle with number
        circ = slide.shapes.add_shape(
            9, Inches(0.8), Inches(y + 0.03), Inches(0.2), Inches(0.2)  # 9 = OVAL
        )
        circ.fill.solid()
        circ.fill.fore_color.rgb = LILLY_RED
        circ.line.fill.background()
        tc = circ.text_frame
        tc.word_wrap = False
        pc = tc.paragraphs[0]
        pc.text = str(i + 1)
        pc.font.size = Pt(8)
        pc.font.bold = True
        pc.font.color.rgb = WHITE
        pc.font.name = "Arial"
        pc.alignment = PP_ALIGN.CENTER

        # Definition text
        _add_text(slide, 1.15, y, 11.2, 0.3, defn, 9, color=DARK)

    # --- Bordered box for SCOPE ---
    scope_y = 4.8
    box2 = slide.shapes.add_shape(5, Inches(0.5), Inches(scope_y), Inches(12.3), Inches(1.5))
    box2.fill.solid()
    box2.fill.fore_color.rgb = WHITE
    box2.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    box2.line.width = Pt(1)

    # Navy header bar
    hdr2 = slide.shapes.add_shape(5, Inches(1.2), Inches(scope_y + 0.15), Inches(3.5), Inches(0.32))
    hdr2.fill.solid()
    hdr2.fill.fore_color.rgb = RGBColor(0x1E, 0x27, 0x61)
    hdr2.line.fill.background()
    tf2 = hdr2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "SCOPE"
    p2.font.size = Pt(10)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.font.name = "Arial"

    # Scope details
    start_raw = config.get("s3_parquet", {}).get("start_date", "2025-01-01")
    end_raw = config.get("s3_parquet", {}).get("end_date", "2025-12-31")
    try:
        from datetime import datetime as dt
        start_fmt = dt.strptime(start_raw, "%Y-%m-%d").strftime("%B %Y")
        end_fmt = dt.strptime(end_raw, "%Y-%m-%d").strftime("%B %Y")
    except:
        start_fmt, end_fmt = start_raw, end_raw

    scope_items = [f"Analysis Period: {start_fmt} – {end_fmt}"]
    country_name = config.get("_country_name", None)
    if country_name:
        scope_items.append(f"Country: {country_name}")
    scope_items.append("This analysis excludes Non-Actionable Insights")

    nb = slide.shapes.add_textbox(Inches(1.0), Inches(scope_y + 0.55), Inches(11), Inches(0.7))
    tf3 = nb.text_frame
    tf3.word_wrap = True
    for j, item in enumerate(scope_items):
        if j == 0:
            p = tf3.paragraphs[0]
        else:
            p = tf3.add_paragraph()
        p.text = f"▸ {item}"
        p.font.size = Pt(10)
        p.font.color.rgb = DARK
        p.font.name = "Arial"


def _slide_thank_you(prs):
    """Final slide: Thank You — uses 'Cover Slide 3' layout."""
    layout = _get_layout(prs, "Cover Slide 3")
    slide = prs.slides.add_slide(layout)

    for ph in slide.placeholders:
        if "Title" in ph.name:
            ph.text = "Thank You"
            for p in ph.text_frame.paragraphs:
                p.font.size = Pt(36)
                p.font.bold = True
                p.font.color.rgb = WHITE
                p.font.name = "Arial"
                p.alignment = PP_ALIGN.CENTER
        elif "Text" in ph.name:
            ph.text = ""


# ============================================================================
# Main
# ============================================================================

def run_report_generator(
    kpi_results: dict[str, pd.DataFrame],
    anomalies: list,
    config: dict,
    output_dir: str = "./data/processed",
) -> Path:
    """Generate the PPTX using the Lilly template."""
    from src.anomaly_detection import anomalies_to_dataframe

    logger.info("=" * 50)
    logger.info("REPORT GENERATOR — START")
    logger.info("=" * 50)

    anomalies_df = anomalies_to_dataframe(anomalies)

    # Load actionable data for executive view
    processed_dir = Path(output_dir)
    parquet_path = processed_dir / config.get("output", {}).get("parquet_file", "actionable_suggestions.parquet")
    df = pd.read_parquet(parquet_path) if parquet_path.exists() else pd.DataFrame()

    # Find template
    template_name = config.get("output", {}).get("template_pptx", "")
    template = Path(template_name) if template_name else None

    if template and not template.exists():
        # Try common locations
        for try_path in [Path("data") / template_name, processed_dir / template_name]:
            if try_path.exists():
                template = try_path
                break

    if template and template.exists():
        logger.info(f"Using template: {template}")
        prs = Presentation(str(template))
        _delete_all_slides(prs)
    else:
        if template_name:
            logger.warning(f"Template not found: {template_name}. Creating without template.")
        else:
            logger.info("No template configured. Creating clean presentation.")
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    # Build slides
    _slide_cover(prs, config)
    _slide_objective(prs, config)
    _slide_executive_view(prs, df, kpi_results)
    _slide_kpi_performance(prs, kpi_results, "4")
    _slide_product_detail(prs, kpi_results, "5")
    _slide_usecase_heatmap(prs, kpi_results, "6")
    next_num = _slide_usecase_performance(prs, kpi_results, start_slide_num=7)
    _slide_anomalies(prs, anomalies_df, str(next_num))
    _slide_appendix_cover(prs)
    _slide_appendix_content(prs, config, str(next_num + 2))
    _slide_thank_you(prs)

    # Save
    country_filter = config.get("_country_filter", None)
    if country_filter:
        filename = f"performance_report_{country_filter}.pptx"
    else:
        filename = "performance_report.pptx"
    out_path = processed_dir / filename
    prs.save(str(out_path))
    logger.info(f"Saved: {out_path} ({out_path.stat().st_size / 1024:.0f} KB)")
    logger.info("=" * 50)
    logger.info(f"REPORT GENERATOR COMPLETE — {len(prs.slides)} slides")
    logger.info("=" * 50)

    return out_path